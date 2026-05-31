"""
文档解析服务 - 支持 PDF/PPTX/DOCX/Markdown
"""
import re
import time
import logging
from pathlib import Path
from typing import List, Tuple, Dict, Optional
import config

logger = logging.getLogger(__name__)


def parse_document(file_path: str) -> Tuple[str, str]:
    """
    解析文档，返回 (纯文本, 文件类型)

    Args:
        file_path: 文件路径

    Returns:
        (提取的文本内容, 文件类型)
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return parse_pdf_document(str(path))

    parsers = {
        ".pptx": parse_pptx,
        ".ppt": parse_pptx,
        ".docx": parse_docx,
        ".doc": parse_docx,
        ".md": parse_markdown,
        ".markdown": parse_markdown,
        ".txt": parse_txt,
    }

    parser = parsers.get(suffix)
    if not parser:
        raise ValueError(f"不支持的文件格式: {suffix}")

    text = parser(str(path))
    return text, suffix.lstrip(".")


def parse_pdf(file_path: str) -> str:
    """解析PDF文件"""
    text, _ = parse_pdf_document(file_path)
    return text


def parse_pdf_document(file_path: str) -> Tuple[str, str]:
    """Parse PDF and return text plus the chunking-friendly content type."""
    if getattr(config, "PDF_PARSER", "pymupdf") == "mineru":
        try:
            mineru_text = _parse_pdf_with_mineru(file_path)
            if mineru_text and mineru_text.strip():
                return mineru_text, "markdown"
        except Exception as e:
            logger.warning("MinerU PDF parsing failed, falling back to PyMuPDF: %s", e)

    return _parse_pdf_with_pymupdf(file_path), "pdf"


def _parse_pdf_with_pymupdf(file_path: str) -> str:
    import fitz  # PyMuPDF
    doc = fitz.open(file_path)
    texts = []
    for page in doc:
        text = page.get_text("text")
        if text.strip():
            texts.append(text.strip())
    doc.close()
    return "\n\n".join(texts)


def _parse_pdf_with_mineru(file_path: str) -> Optional[str]:
    import requests

    base_url = getattr(config, "MINERU_API_BASE", "").rstrip("/")
    if not base_url:
        return None

    path = Path(file_path)
    headers = _mineru_headers()
    timeout = int(getattr(config, "MINERU_TIMEOUT", 300))

    create_resp = requests.post(
        f"{base_url}/parse/file",
        json=_mineru_create_payload(path),
        headers=headers,
        timeout=timeout,
    )
    create_resp.raise_for_status()
    create_data = _mineru_data(create_resp.json())
    task_id = create_data.get("task_id") or create_data.get("taskId") or create_data.get("id")
    upload_url = (
        create_data.get("upload_url")
        or create_data.get("uploadUrl")
        or create_data.get("file_url")
        or create_data.get("fileUrl")
    )
    if not task_id or not upload_url:
        return None

    with open(file_path, "rb") as pdf:
        upload_resp = requests.put(upload_url, data=pdf, timeout=timeout)
    upload_resp.raise_for_status()

    result = _wait_for_mineru_result(base_url, task_id, headers)
    return _download_mineru_markdown(result, headers)


def _mineru_create_payload(path: Path) -> Dict[str, object]:
    return {
        "filename": path.name,
        "file_name": path.name,
        "language": getattr(config, "MINERU_LANGUAGE", "ch"),
        "enable_table": bool(getattr(config, "MINERU_ENABLE_TABLE", True)),
        "enable_formula": bool(getattr(config, "MINERU_ENABLE_FORMULA", True)),
        "is_ocr": bool(getattr(config, "MINERU_IS_OCR", False)),
    }


def _mineru_headers() -> Dict[str, str]:
    headers = {"Content-Type": "application/json"}
    token = getattr(config, "MINERU_API_TOKEN", "")
    if token:
        headers["Authorization"] = f"Bearer {token}"
    return headers


def _wait_for_mineru_result(base_url: str, task_id: str, headers: Dict[str, str]) -> Dict[str, object]:
    import requests

    timeout = int(getattr(config, "MINERU_TIMEOUT", 300))
    interval = max(1, int(getattr(config, "MINERU_POLL_INTERVAL", 3)))
    deadline = time.time() + timeout

    while time.time() < deadline:
        resp = requests.get(f"{base_url}/parse/{task_id}", headers=headers, timeout=timeout)
        resp.raise_for_status()
        data = _mineru_data(resp.json())
        status = str(data.get("status") or data.get("state") or "").lower()
        if status in {"done", "finished", "success", "succeeded", "completed"}:
            return data
        if status in {"failed", "fail", "error"}:
            return {}
        time.sleep(interval)

    return {}


def _download_mineru_markdown(result: Dict[str, object], headers: Dict[str, str]) -> Optional[str]:
    import requests

    markdown = result.get("markdown") or result.get("md") or result.get("content")
    if isinstance(markdown, str) and markdown.strip():
        return markdown

    md_url = (
        result.get("markdown_url")
        or result.get("markdownUrl")
        or result.get("md_url")
        or result.get("mdUrl")
    )
    if not md_url:
        return None

    resp = requests.get(md_url, headers=headers, timeout=int(getattr(config, "MINERU_TIMEOUT", 300)))
    resp.raise_for_status()
    return resp.text


def _mineru_data(payload: Dict[str, object]) -> Dict[str, object]:
    data = payload.get("data", payload) if isinstance(payload, dict) else {}
    return data if isinstance(data, dict) else {}


def parse_pptx(file_path: str) -> str:
    """解析PPT/PPTX文件"""
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
            elif shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        slide_texts.append(row_text)
        if slide_texts:
            texts.append(f"[幻灯片 {slide_num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)


def parse_docx(file_path: str) -> str:
    """解析DOCX文件"""
    from docx import Document
    doc = Document(file_path)
    texts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            texts.append(text)
    # 也提取表格内容
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                texts.append(row_text)
    return "\n\n".join(texts)


def parse_markdown(file_path: str) -> str:
    """解析Markdown文件"""
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def parse_txt(file_path: str) -> str:
    """解析纯文本文件"""
    # 尝试多种编码
    for encoding in ["utf-8", "gbk", "gb2312", "utf-16"]:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                return f.read()
        except (UnicodeDecodeError, UnicodeError):
            continue
    raise ValueError(f"无法解码文件: {file_path}")


def chunk_text(text: str, chunk_size: int = None, overlap: int = None) -> List[str]:
    """
    将长文本分割为固定大小的片段

    Args:
        text: 原始文本
        chunk_size: 片段大小（字符数）
        overlap: 片段重叠大小

    Returns:
        文本片段列表
    """
    chunk_size = chunk_size or config.CHUNK_SIZE
    overlap = overlap or config.CHUNK_OVERLAP

    if not text.strip():
        return []

    # 按段落分割，保持语义完整性
    paragraphs = text.split("\n\n")
    chunks = []
    current_chunk = ""

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        if len(current_chunk) + len(para) + 2 <= chunk_size:
            current_chunk = current_chunk + "\n\n" + para if current_chunk else para
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            # 如果单个段落就超长，强制分割
            if len(para) > chunk_size:
                for i in range(0, len(para), chunk_size - overlap):
                    sub = para[i:i + chunk_size]
                    if sub.strip():
                        chunks.append(sub.strip())
                current_chunk = ""
            else:
                current_chunk = para

    if current_chunk.strip():
        chunks.append(current_chunk.strip())

    return chunks


def chunk_markdown_by_headings(text: str, chunk_size: int = None, overlap: int = None) -> List[Dict[str, str]]:
    """Split Markdown by heading sections and keep heading paths as metadata."""
    chunk_size = chunk_size or config.CHUNK_SIZE
    overlap = overlap or config.CHUNK_OVERLAP

    if not text.strip():
        return []

    sections = _split_markdown_sections(text)
    if not sections:
        return [
            {"content": chunk, "title_path": ""}
            for chunk in chunk_text(text, chunk_size=chunk_size, overlap=overlap)
        ]

    chunks = []
    for section in sections:
        chunks.extend(_chunk_markdown_section(section, chunk_size, overlap))
    return [chunk for chunk in chunks if chunk.get("content", "").strip()]


def _split_markdown_sections(text: str) -> List[Dict[str, str]]:
    heading_re = re.compile(r"^(#{1,6})\s+(.+?)\s*#*\s*$")
    sections = []
    current = None
    heading_stack = []
    preamble = []
    in_code_fence = False

    for line in text.splitlines():
        stripped = line.strip()
        if stripped.startswith("```") or stripped.startswith("~~~"):
            in_code_fence = not in_code_fence

        match = heading_re.match(line) if not in_code_fence else None
        if match:
            if current:
                sections.append(current)
            elif preamble and "".join(preamble).strip():
                sections.append({
                    "path": "Preamble",
                    "text": "\n".join(preamble).strip(),
                })
                preamble = []

            level = len(match.group(1))
            title = _clean_markdown_heading(match.group(2))
            heading_stack = heading_stack[:level - 1]
            heading_stack.append(title)
            current = {
                "path": " > ".join([h for h in heading_stack if h]),
                "text": "",
            }
            continue

        if current:
            current["text"] = current["text"] + "\n" + line if current["text"] else line
        else:
            preamble.append(line)

    if current:
        sections.append(current)
    elif preamble and "".join(preamble).strip():
        sections.append({
            "path": "Preamble",
            "text": "\n".join(preamble).strip(),
        })

    return sections if any(section["path"] != "Preamble" for section in sections) else []


def _chunk_markdown_section(section: Dict[str, str], chunk_size: int, overlap: int) -> List[Dict[str, str]]:
    path = section.get("path") or "Untitled"
    section_text = (section.get("text") or "").strip()

    if not section_text:
        return []

    if len(section_text) <= chunk_size:
        return [{"content": section_text, "title_path": path}]

    section_chunks = chunk_text(section_text, chunk_size=chunk_size, overlap=overlap)
    return [{"content": chunk, "title_path": path} for chunk in section_chunks]


def _clean_markdown_heading(title: str) -> str:
    title = re.sub(r"\s+", " ", title or "").strip()
    return title.strip("#").strip()
